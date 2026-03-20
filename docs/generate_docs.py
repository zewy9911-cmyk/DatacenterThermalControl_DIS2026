"""
generate_docs.py  —  Creates prezentaciya.pptx and dokumentaciya.docx
Run:  python docs/generate_docs.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DI, Pt as DP, RGBColor as DR
from docx.enum.text import WD_ALIGN_PARAGRAPH

OUT = os.path.dirname(os.path.abspath(__file__))

# ── colours ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x39, 0x64)
BLUE   = RGBColor(0x2E, 0x75, 0xB6)
LBLUE  = RGBColor(0x9D, 0xC3, 0xE6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x20, 0x20, 0x20)
GRAY   = RGBColor(0x60, 0x60, 0x60)
GREEN  = RGBColor(0x37, 0x86, 0x34)

# ════════════════════════════════════════════════════════════════════════════
#  POWERPOINT
# ════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK   = prs.slide_layouts[6]
TITLED  = prs.slide_layouts[1]   # title + content

def bg(slide, rgb):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb

def txb(slide, text, l, t, w, h, sz=18, bold=False, col=BLACK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text       = text
    r.font.size  = Pt(sz)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = col
    r.font.name  = "Calibri"

def header_bar(slide, title_text):
    """Blue header rectangle + white title text"""
    bar = slide.shapes.add_shape(1,
        Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    txb(slide, title_text, 0.25, 0.1, 12.8, 0.9,
        sz=24, bold=True, col=WHITE, align=PP_ALIGN.LEFT)

def footer(slide, page):
    txb(slide, f"Разпределени Вградени Системи  |  Курсов проект  |  2025/2026  |  {page}",
        0.3, 7.15, 12.7, 0.3, sz=9, col=GRAY, align=PP_ALIGN.CENTER)
    foot_bar = slide.shapes.add_shape(1,
        Inches(0), Inches(7.35), Inches(13.33), Inches(0.15))
    foot_bar.fill.solid(); foot_bar.fill.fore_color.rgb = BLUE
    foot_bar.line.fill.background()

def bullets(slide, items, l, t, w, h, sz=17, col=BLACK, indent="  • "):
    tb  = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{indent}{item}"
        r.font.size = Pt(sz)
        r.font.color.rgb = col
        r.font.name = "Calibri"

def small_table(slide, headers, rows, l, t, w):
    cols = len(headers)
    tbl  = slide.shapes.add_table(len(rows)+1, cols,
                Inches(l), Inches(t),
                Inches(w), Inches(0.45*(len(rows)+1))).table
    tbl.first_row = True
    # header row
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.bold  = True
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.size  = Pt(13)
        p.runs[0].font.name  = "Calibri"
    # data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.text = val
            if ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE7,0xF0,0xFA)
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(12)
            p.runs[0].font.name = "Calibri"
            p.runs[0].font.color.rgb = BLACK

# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
# top accent
bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.18))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
# logo placeholder text
txb(s, "🏫  Технически Университет", 0.4, 0.25, 12, 0.5, sz=13, col=LBLUE, align=PP_ALIGN.LEFT)
txb(s, "Разпределена Вградена Система за\nТермично Управление на Център за Данни",
    0.8, 1.1, 11.7, 2.3, sz=34, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
txb(s, "Курсов проект", 0.8, 3.45, 11.7, 0.55, sz=22, col=LBLUE, align=PP_ALIGN.CENTER)
txb(s, "Дисциплина:  Разпределени Вградени Системи",
    0.8, 4.05, 11.7, 0.5, sz=15, col=LBLUE, align=PP_ALIGN.CENTER)
# divider
div = s.shapes.add_shape(1, Inches(3.5), Inches(4.65), Inches(6.3), Inches(0.05))
div.fill.solid(); div.fill.fore_color.rgb = BLUE; div.line.fill.background()
# names block
txb(s, "Изготвили:", 0.8, 4.8, 11.7, 0.4, sz=13, col=LBLUE, align=PP_ALIGN.CENTER, italic=True)
txb(s, "Стефани Узунова  —  ФН: 616551\n"
       "Владимир Върбанов  —  ФН: 616602\n"
       "Данаил Атанасов  —  ФН: 616614",
    0.8, 5.2, 11.7, 1.1, sz=16, col=WHITE, align=PP_ALIGN.CENTER)
txb(s, "Март  2026 г.", 0.8, 6.4, 11.7, 0.4, sz=13, col=LBLUE, align=PP_ALIGN.CENTER)
# bottom bar
b = s.shapes.add_shape(1, Inches(0), Inches(7.3), Inches(13.33), Inches(0.2))
b.fill.solid(); b.fill.fore_color.rgb = BLUE; b.line.fill.background()

# ── SLIDE 2 — Съдържание ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header_bar(s, "Съдържание")
items = [
    "1.  Описание на задачата",
    "2.  Системен обзор и хардуерни компоненти",
    "3.  GPIO разпределение на пиновете",
    "4.  Модул 1 — Следене на температурата и влажността",
    "5.  Модул 2 — Управление на вентилацията (вентилатор + клапи)",
    "6.  Логика на автоматично управление",
    "7.  Комуникационни протоколи — MQTT и SNMP",
    "8.  Програмна архитектура (Python / FastAPI / React)",
    "9.  Уеб интерфейс и симулационен режим",
    "10. Заключение и резултати",
]
bullets(s, items, 0.6, 1.25, 12, 5.8, sz=18, indent="")
footer(s, "2 / 13")

# ── SLIDE 3 — Задача ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header_bar(s, "Описание на задачата")
txb(s, "Задание:", 0.6, 1.2, 12, 0.4, sz=17, bold=True, col=NAVY)
bullets(s, [
    "Проектиране на разпределена вградена система за термично управление на ЦД",
    "Реализация на Raspberry Pi 5 с Python FastAPI бекенд и React уеб интерфейс",
], 0.6, 1.6, 12, 0.9, sz=16)
txb(s, "Модул 1 — Следене на температурата:", 0.6, 2.55, 12, 0.4, sz=17, bold=True, col=NAVY)
bullets(s, [
    "Отчитане на температура и влажност чрез сензор BME280 (I²C) или DHT22 (GPIO)",
    "Докладване на стойностите по мрежа чрез SNMP v2c протокол",
    "Съхранение на исторически данни в SQLite база данни",
], 0.6, 2.95, 12, 1.1, sz=16)
txb(s, "Модул 2 — Управление на вентилацията:", 0.6, 4.1, 12, 0.4, sz=17, bold=True, col=NAVY)
bullets(s, [
    "Управление на вентилатор с 3 степени (НИСКА / СРЕДНА / ВИСОКА) чрез Modbus RTU",
    "Управление на 2 клапи: Клапа 1 — рециркулация вътре, Клапа 2 — изпускане навън",
    "Детекция на въздушен поток чрез дискретни входове (DI) на Modbus модула",
    "Автоматичен и ръчен режим на управление + MQTT брокер за много Pi устройства",
], 0.6, 4.5, 12, 1.4, sz=16)
footer(s, "3 / 13")

# ── SLIDE 4 — Хардуер ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header_bar(s, "Хардуерни компоненти")
small_table(s,
    ["Компонент", "Модел / Спецификация", "Интерфейс", "Брой"],
    [
        ["Едноплаткови компютър", "Raspberry Pi 5  (4 GB RAM)", "—", "1"],
        ["Сензор температура/влажност", "BME280", "I²C  0x76", "1"],
        ["Алт. сензор", "DHT22", "GPIO 4  (BCM)", "1"],
        ["RS-485 USB адаптер", "CH340 / FT232", "USB → /dev/ttyUSB0", "1"],
        ["Modbus RTU релеен модул", "8-канален, 5A/230VAC", "RS-485  Slave ID 1", "1"],
        ["Вентилатор", "3-степенен AC/DC", "Relay 1–3  (one-hot)", "1"],
        ["Актуатори за клапи", "24 VDC моторизирани", "Relay 4–5", "2"],
        ["Датчици за въздушен поток", "NPN/PNP дигитален изход", "Modbus DI 0, DI 1", "2"],
    ], 0.35, 1.2, 12.6)
footer(s, "4 / 13")

# ── SLIDE 5 — GPIO ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header_bar(s, "GPIO разпределение — Raspberry Pi 5 (BCM номерация)")
small_table(s,
    ["GPIO (BCM)", "Физ. пин", "Посока", "Функция"],
    [
        ["GPIO  2 (SDA)", "Pin 3", "Изход/вход", "BME280 — I²C Data"],
        ["GPIO  3 (SCL)", "Pin 5", "Изход",      "BME280 — I²C Clock"],
        ["GPIO  4",       "Pin 7", "Вход",        "DHT22 — данни (алт. сензор)"],
        ["GPIO 17",       "Pin 11","Изход",        "Статус LED (зелен — heartbeat)"],
        ["GPIO 27",       "Pin 13","Изход",        "Алarm LED (червен — критична T°)"],
        ["GPIO 22",       "Pin 15","Вход (PU)",    "Бутон — превключване АВТО/РЪЧЕН"],
        ["GPIO 23",       "Pin 16","Вход",         "Тахометър на вентилатора (NPN)"],
        ["GPIO 24",       "Pin 18","Вход (PD)",    "Обратна връзка — Клапа 1"],
        ["GPIO 25",       "Pin 22","Вход (PD)",    "Обратна връзка — Клапа 2"],
    ], 0.35, 1.2, 12.6)
txb(s, "PU = Pull-Up  |  PD = Pull-Down  |  Modbus DI-0/DI-1 се четат директно от RS-485 модула",
    0.4, 6.65, 12.5, 0.4, sz=11, col=GRAY, italic=True)
footer(s, "5 / 13")

# ── SLIDE 6 — Модул 1 ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header_bar(s, "Модул 1 — Следене на температурата и влажността")
txb(s, "Сензор BME280 / DHT22", 0.6, 1.2, 6, 0.4, sz=17, bold=True, col=NAVY)
bullets(s, [
    "BME280 — I²C, адрес 0x76, ±0.5 °C точност",
    "DHT22  — GPIO 4, алтернатива при липса на I²C",
    "Автоматичен fallback към симулация при грешка",
    "Четене на всеки 5 секунди (конфигурируемо)",
], 0.6, 1.6, 6.2, 1.6, sz=15)
txb(s, "SNMP v2c докладване", 0.6, 3.3, 6, 0.4, sz=17, bold=True, col=NAVY)
bullets(s, [
    "Enterprise OID:  1.3.6.1.4.1.54321",
    "Температура → .1.1.0  (×10, напр. 264 = 26.4 °C)",
    "Влажност    → .1.2.0  (×10)",
    "Автоматични трапове при T ≥ 30 °C (.4.1)",
    "Автоматични трапове при RH ≥ 70 % (.4.2)",
    "Community: public  |  UDP порт 161/162",
], 0.6, 3.7, 6.2, 2.1, sz=15)
txb(s, "SQLite база данни", 7.2, 1.2, 5.8, 0.4, sz=17, bold=True, col=NAVY)
bullets(s, [
    "Таблица readings — запис на 60 секунди",
    "Таблица events   — системни събития",
    "API: GET /api/history?hours=24",
    "API: GET /api/statistics  (min/max/avg)",
    "aiosqlite — асинхронен достъп",
], 7.2, 1.6, 5.8, 2.0, sz=15)
footer(s, "6 / 13")

# ── SLIDE 7 — Модул 2 ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header_bar(s, "Модул 2 — Управление на вентилацията")
txb(s, "Modbus RTU — релейни изходи", 0.6, 1.2, 6, 0.4, sz=17, bold=True, col=NAVY)
small_table(s,
    ["Relay", "Coil", "Функция"],
    [
        ["Relay 1", "Coil 0", "Вентилатор НИСКА (~33 %)"],
        ["Relay 2", "Coil 1", "Вентилатор СРЕДНА (~66 %)"],
        ["Relay 3", "Coil 2", "Вентилатор ВИСОКА (100 %)"],
        ["Relay 4", "Coil 3", "Клапа 1 — Рециркулация"],
        ["Relay 5", "Coil 4", "Клапа 2 — Изпускане навън"],
    ], 0.5, 1.7, 6.0)
txb(s, "Клапи — режими на работа", 7.0, 1.2, 6.0, 0.4, sz=17, bold=True, col=NAVY)
bullets(s, [
    "♻  Рециркулация (T < 25 °C):",
    "     Клапа 1 ОТВОРЕНА + Клапа 2 ЗАТВОРЕНА",
    "     → въздухът се върти вътре в стаята",
    "",
    "🌬  Изпускане (T ≥ 25 °C):",
    "     Клапа 1 ЗАТВОРЕНА + Клапа 2 ОТВОРЕНА",
    "     → горещ въздух се изкарва навън",
    "",
    "⚠  Fault override: ако DI-1 (изходящ поток)",
    "     не е засечен → автоматично рециркулация",
], 7.0, 1.6, 6.0, 4.8, sz=14, indent="")
footer(s, "7 / 13")

# ── SLIDE 8 — Логика на управление ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header_bar(s, "Логика на автоматично управление")
small_table(s,
    ["Температура", "Скорост вентилатор", "Клапа 1 (Рецирк.)", "Клапа 2 (Изпуск.)"],
    [
        ["T < 20 °C",       "⭕  ИЗКЛЮЧЕН",  "✅  ОТВОРЕНА", "❌  ЗАТВОРЕНА"],
        ["20 ≤ T < 25 °C",  "🟢  НИСКА",    "✅  ОТВОРЕНА", "❌  ЗАТВОРЕНА"],
        ["25 ≤ T < 30 °C",  "🟡  СРЕДНА",   "❌  ЗАТВОРЕНА","✅  ОТВОРЕНА"],
        ["T ≥ 30 °C",       "🔴  ВИСОКА",   "❌  ЗАТВОРЕНА","✅  ОТВОРЕНА"],
    ], 0.5, 1.2, 12.3)
txb(s, "Режими на управление:", 0.6, 4.0, 12, 0.4, sz=17, bold=True, col=NAVY)
bullets(s, [
    "АВТО — системата автоматично избира вентилатор и клапи по температура",
    "РЪЧЕН — операторът задава вентилатор (0–3) и всяка клапа независимо",
    "Превключване: бутон GPIO 22  или  REST API  POST /api/mode  или  MQTT команда",
    "Прагове конфигурируемо чрез API:  POST /api/thresholds  {\"off\": 20, \"low\": 25, \"medium\": 30}",
], 0.6, 4.4, 12, 2.2, sz=16)
footer(s, "8 / 13")

# ── SLIDE 9 — MQTT & SNMP ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header_bar(s, "Комуникационни протоколи — MQTT и SNMP")
txb(s, "MQTT (paho-mqtt) — Mosquitto брокер", 0.6, 1.2, 6.2, 0.4, sz=16, bold=True, col=NAVY)
bullets(s, [
    "Broker: TCP :1883  |  TLS :8883 (опция)",
    "NODE_ID: datacenter_node_01  (уникален за Pi)",
    "Публикуване: datacenter/{node}/status (retained JSON)",
    "Индивидуални топики: sensors/temperature, actuators/fan_speed...",
    "Команди (subscribe): cmd/fan_speed, cmd/valve, cmd/mode",
    "Fallback: при > 10 неуспешни опита → SNMP режим",
    "LWT (Last Will): status = offline при изключване",
], 0.6, 1.65, 6.2, 3.1, sz=14)
txb(s, "SNMP v2c — резервен протокол", 7.0, 1.2, 6.0, 0.4, sz=16, bold=True, col=NAVY)
small_table(s,
    ["OID суфикс", "Описание"],
    [
        [".1.1.0", "Температура (×10 °C)"],
        [".1.2.0", "Влажност (×10 %)"],
        [".2.1.0", "Скорост вентилатор (0–3)"],
        [".2.2.0", "Клапа 1 (0/1)"],
        [".2.3.0", "Клапа 2 (0/1)"],
        [".3.1.0", "Въздух вход (0/1)"],
        [".3.2.0", "Въздух изход (0/1)"],
    ], 6.9, 1.65, 6.1)
footer(s, "9 / 13")

# ── SLIDE 10 — Програмна архитектура ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header_bar(s, "Програмна архитектура")
txb(s, "Бекенд (Python 3.11):", 0.6, 1.2, 6.2, 0.4, sz=16, bold=True, col=NAVY)
small_table(s,
    ["Файл", "Описание"],
    [
        ["main.py",             "FastAPI сървър — REST API, WebSocket, статични файлове"],
        ["config.py",           "Всички константи — GPIO, Modbus, MQTT, SNMP, прагове"],
        ["sensor_reader.py",    "BME280/DHT22 + симулационен override"],
        ["control_logic.py",    "Автоматична логика АВТО/РЪЧЕН режим"],
        ["IOControl/modbus_control.py", "Modbus RTU — вентилатор + клапи + DI"],
        ["IOControl/gpio_handler.py",   "GPIO — LED, бутон, тахометър, feedback"],
        ["Networking/mqtt_client.py",   "MQTT pub/sub + fallback логика"],
        ["Networking/snmp_agent.py",    "SNMP v2c агент + traps"],
        ["Storage/data_logger.py",      "SQLite асинхронен логер (aiosqlite)"],
    ], 0.4, 1.65, 7.5)
txb(s, "Фронтенд (React 18 + Vite):", 8.2, 1.2, 5.0, 0.4, sz=16, bold=True, col=NAVY)
bullets(s, [
    "SensorGauges.jsx — радиални циферблати",
    "FanControl.jsx — 4 бутона (OFF/LOW/MED/HIGH)",
    "ValveControl.jsx — 2 toggle превключвателя",
    "AirflowStatus.jsx — DI индикатори",
    "ModeToggle.jsx — АВТО/РЪЧЕН",
    "AlertsPanel.jsx — цветни известия",
    "HistoricalChart.jsx — Recharts линейна графика",
    "SimulationPanel.jsx — слайдери + сценарии",
], 8.2, 1.65, 4.8, 4.5, sz=13)
footer(s, "10 / 13")

# ── SLIDE 11 — Уеб интерфейс ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header_bar(s, "Уеб интерфейс — React Dashboard")
txb(s, "Достъп:  http://<RPi-IP>:8000   |   API документация:  /docs   |   WebSocket:  ws://<IP>:8000/ws",
    0.5, 1.2, 12.3, 0.45, sz=14, col=NAVY, bold=True)
txb(s, "Функционалности:", 0.6, 1.75, 12, 0.4, sz=16, bold=True, col=NAVY)
bullets(s, [
    "Радиални циферблати за температура (°C) и влажност (%) с цветово кодиране",
    "Панел АВТО/РЪЧЕН режим — превключване с един клик",
    "Управление на вентилатора — 4 бутона (OFF / НИСКА / СРЕДНА / ВИСОКА)",
    "Управление на клапите — toggle превключватели (само в РЪЧЕН режим)",
    "Датчици за въздушен поток — визуален статус на DI-0 и DI-1",
    "GPIO обратна връзка — RPM на вентилатора, позиция на клапите",
    "Известия — цветово кодирани (зелен/жълт/червен)",
    "Историческа графика — линейна Recharts за последните 1/6/12/24 часа",
    "Статус на MQTT — ✓ свързан / SNMP↩ fallback",
], 0.6, 2.15, 8.5, 4.5, sz=15)
txb(s, "REST API:", 9.3, 1.75, 3.8, 0.4, sz=16, bold=True, col=NAVY)
bullets(s, [
    "GET   /api/status",
    "POST  /api/fan",
    "POST  /api/valve",
    "POST  /api/mode",
    "POST  /api/thresholds",
    "GET   /api/history",
    "GET   /api/statistics",
    "GET   /api/mqtt",
    "GET   /api/gpio",
    "GET/POST /api/simulation",
], 9.3, 2.15, 3.8, 4.5, sz=13)
footer(s, "11 / 13")

# ── SLIDE 12 — Симулация ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header_bar(s, "Симулационен режим")
txb(s, "Предназначение:", 0.6, 1.2, 12, 0.4, sz=16, bold=True, col=NAVY)
bullets(s, [
    "Системата може да работи напълно без физически хардуер (сензори / Modbus)",
    "Два нива на симулация — автоматична (няма хардуер) и ръчна (API override)",
], 0.6, 1.6, 12, 0.7, sz=15)
txb(s, "Ниво 1 — Автоматична симулация (hardware fallback):", 0.6, 2.4, 12, 0.4, sz=16, bold=True, col=NAVY)
bullets(s, [
    "При липса на BME280/DHT22 — sensor_reader.py генерира реалистични стойности (синусоидален drift)",
    "При липса на Modbus адаптер — IOControl симулира relay/DI операции с лог съобщения",
], 0.6, 2.8, 12, 0.75, sz=15)
txb(s, "Ниво 2 — Ръчен override (SimulationPanel в UI):", 0.6, 3.65, 12, 0.4, sz=16, bold=True, col=NAVY)
bullets(s, [
    "Слайдер за температура (10–45 °C) и влажност (10–95 %) в реално време",
    "Предварителен преглед на очакваното поведение (вентилатор + клапи) преди прилагане",
    "Бързи сценарии: 🌿 Нормална / ☀ Топло / 🌡 Горещо / 🔥 Критично / 💧 Вис. влажност",
    "API: POST /api/simulation  {\"enabled\": true, \"temperature\": 31, \"humidity\": 68}",
    "Индикатор ⚗ СИМУЛАЦИЯ се показва в заглавието на сензорния панел",
], 0.6, 4.05, 12, 2.0, sz=15)
txb(s, "Полезно за:  демонстрация без хардуер  |  тестване на UI/UX  |  верификация на логиката",
    0.6, 6.25, 12.2, 0.45, sz=13, col=GRAY, italic=True)
footer(s, "12 / 13")

# ── SLIDE 13 — Заключение ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header_bar(s, "Заключение и резултати")
txb(s, "Реализирани функционалности:", 0.6, 1.2, 12, 0.4, sz=16, bold=True, col=NAVY)
small_table(s,
    ["Функционалност", "Статус"],
    [
        ["Четене на температура/влажност (BME280 / DHT22)", "✅  Реализирано"],
        ["Modbus RTU управление — вентилатор 3 степени + 2 клапи", "✅  Реализирано"],
        ["GPIO — LED индикатори, бутон, тахометър, valve feedback", "✅  Реализирано"],
        ["SNMP v2c докладване + автоматични трапове", "✅  Реализирано"],
        ["MQTT брокер (Mosquitto) + multi-node поддръжка", "✅  Реализирано"],
        ["Автоматичен fallback MQTT → SNMP", "✅  Реализирано"],
        ["SQLite база данни с история и статистика", "✅  Реализирано"],
        ["React уеб интерфейс с WebSocket в реално време", "✅  Реализирано"],
        ["Симулационен режим (без хардуер)", "✅  Реализирано"],
    ], 0.4, 1.65, 12.5)
txb(s, "Благодарим за вниманието!",
    0.5, 6.3, 12.3, 0.55, sz=22, bold=True, col=NAVY, align=PP_ALIGN.CENTER)
footer(s, "13 / 13")

prs.save(os.path.join(OUT, "prezentaciya.pptx"))
print("✓ prezentaciya.pptx saved")

# ════════════════════════════════════════════════════════════════════════════
#  WORD DOCUMENT
# ════════════════════════════════════════════════════════════════════════════

doc = Document()

# Page margins
for section in doc.sections:
    section.top_margin    = DI(1)
    section.bottom_margin = DI(1)
    section.left_margin   = DI(1.2)
    section.right_margin  = DI(1.2)

def dh1(text):
    p = doc.add_heading(text, level=1)
    p.runs[0].font.color.rgb = DR(0x1F, 0x39, 0x64)

def dh2(text):
    p = doc.add_heading(text, level=2)
    p.runs[0].font.color.rgb = DR(0x2E, 0x75, 0xB6)

def dh3(text):
    p = doc.add_heading(text, level=3)

def dp(text, bold=False, italic=False):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold   = bold
    r.italic = italic

def dtable(headers, rows):
    t = doc.add_table(rows=1+len(rows), cols=len(headers))
    t.style = "Table Grid"
    for i, h in enumerate(headers):
        c = t.rows[0].cells[i]
        c.text = h
        c.paragraphs[0].runs[0].bold = True
        c.paragraphs[0].runs[0].font.color.rgb = DR(0xFF,0xFF,0xFF)
        tc = c._tc; tcPr = tc.get_or_add_tcPr()
        from docx.oxml import OxmlElement as OE
        shd = OE('w:shd')
        shd.set('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill','1F3964')
        shd.set('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}color','auto')
        shd.set('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val','clear')
        tcPr.append(shd)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            t.rows[ri+1].cells[ci].text = val
    doc.add_paragraph()

def dbullets(items):
    for item in items:
        doc.add_paragraph(item, style="List Bullet")

# ── Title page ────────────────────────────────────────────────────────────────
doc.add_paragraph()
title_p = doc.add_paragraph()
title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = title_p.add_run("ТЕХНИЧЕСКИ УНИВЕРСИТЕТ")
r.bold = True; r.font.size = DP(16)
r.font.color.rgb = DR(0x1F,0x39,0x64)

doc.add_paragraph()
title_p2 = doc.add_paragraph()
title_p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
r2 = title_p2.add_run("КУРСОВ ПРОЕКТ")
r2.bold = True; r2.font.size = DP(20)
r2.font.color.rgb = DR(0x2E,0x75,0xB6)

doc.add_paragraph()
t3 = doc.add_paragraph()
t3.alignment = WD_ALIGN_PARAGRAPH.CENTER
r3 = t3.add_run("Разпределена Вградена Система за Термично Управление на Център за Данни")
r3.bold = True; r3.font.size = DP(18)
r3.font.color.rgb = DR(0x1F,0x39,0x64)

doc.add_paragraph()
t4 = doc.add_paragraph()
t4.alignment = WD_ALIGN_PARAGRAPH.CENTER
r4 = t4.add_run("Дисциплина: Разпределени Вградени Системи")
r4.font.size = DP(14)

doc.add_paragraph()
doc.add_paragraph()
t5 = doc.add_paragraph()
t5.alignment = WD_ALIGN_PARAGRAPH.CENTER
r5 = t5.add_run("Изготвили:")
r5.bold = True; r5.font.size = DP(13)

for name, fn in [("Стефани Узунова","616551"),
                  ("Владимир Върбанов","616602"),
                  ("Данаил Атанасов","616614")]:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run(f"{name}  —  ФН: {fn}").font.size = DP(13)

doc.add_paragraph()
t6 = doc.add_paragraph()
t6.alignment = WD_ALIGN_PARAGRAPH.CENTER
t6.add_run("Март  2026 г.").font.size = DP(13)

doc.add_page_break()

# ── 1. Въведение ─────────────────────────────────────────────────────────────
dh1("1. Въведение")
dp("Настоящият курсов проект разработва разпределена вградена система за термично управление на сървърна стая (Datacenter). Системата е реализирана на платформата Raspberry Pi 5 и включва два основни функционални модула:")
dbullets([
    "Модул 1 — Следене на температурата и влажността с SNMP докладване",
    "Модул 2 — Автоматично управление на вентилацията (вентилатор с 3 степени + 2 клапи)",
])
dp("Комуникацията между множество Pi устройства се осъществява чрез MQTT брокер (Mosquitto), с автоматичен fallback към SNMP при недостъпен брокер.")

# ── 2. Хардуер ───────────────────────────────────────────────────────────────
dh1("2. Хардуерни компоненти")
dh2("2.1. Списък с компоненти")
dtable(
    ["Компонент","Модел","Интерфейс","Бр."],
    [
        ["Едноплаткови компютър","Raspberry Pi 5 (4 GB)","—","1"],
        ["Температурен сензор","BME280","I²C 0x76","1"],
        ["Алт. сензор","DHT22","GPIO 4","1"],
        ["RS-485 адаптер","CH340/FT232","USB /dev/ttyUSB0","1"],
        ["Modbus RTU релеен модул","8-кан. 5A/230VAC","RS-485 Slave 1","1"],
        ["Вентилатор","3-степенен","Relay 1–3","1"],
        ["Клапни актуатори","24 VDC","Relay 4–5","2"],
        ["Въздушни датчици","NPN/PNP DI","Modbus DI 0,1","2"],
    ]
)

dh2("2.2. GPIO разпределение на пиновете (BCM)")
dtable(
    ["GPIO (BCM)","Физ. пин","Посока","Функция"],
    [
        ["GPIO 2 (SDA)","Pin 3","Двупосочен","BME280 I²C Data"],
        ["GPIO 3 (SCL)","Pin 5","Изход","BME280 I²C Clock"],
        ["GPIO 4","Pin 7","Вход","DHT22 данни"],
        ["GPIO 17","Pin 11","Изход","Статус LED (зелен, heartbeat)"],
        ["GPIO 27","Pin 13","Изход","Alarm LED (червен, критична T°)"],
        ["GPIO 22","Pin 15","Вход PU","Бутон АВТО/РЪЧЕН (дебаунс 200 ms)"],
        ["GPIO 23","Pin 16","Вход","Тахометър вентилатор (NPN импулс)"],
        ["GPIO 24","Pin 18","Вход PD","Обратна връзка Клапа 1"],
        ["GPIO 25","Pin 22","Вход PD","Обратна връзка Клапа 2"],
    ]
)

dh2("2.3. Modbus RTU — адресна карта")
dtable(
    ["Тип","Адрес","Функция"],
    [
        ["Coil","0","Вентилатор НИСКА (~33 %)"],
        ["Coil","1","Вентилатор СРЕДНА (~66 %)"],
        ["Coil","2","Вентилатор ВИСОКА (100 %)"],
        ["Coil","3","Клапа 1 — Рециркулация"],
        ["Coil","4","Клапа 2 — Изпускане навън"],
        ["DI","0","Датчик входящ въздушен поток"],
        ["DI","1","Датчик изходящ въздушен поток"],
    ]
)

# ── 3. Програмна архитектура ──────────────────────────────────────────────────
dh1("3. Програмна архитектура")
dh2("3.1. Бекенд модули (Python 3.11)")
dtable(
    ["Файл / Модул","Описание"],
    [
        ["main.py","FastAPI + Uvicorn сървър — REST API, WebSocket, статични файлове"],
        ["config.py","Централна конфигурация: GPIO пинове, Modbus, MQTT, SNMP, прагове"],
        ["sensor_reader.py","Четене от BME280/DHT22; симулационен override (set_override)"],
        ["control_logic.py","Автоматична логика; АВТО/РЪЧЕН режим; управление на прагове"],
        ["IOControl/modbus_control.py","Modbus RTU клиент — fan one-hot + valve coils + DI четене"],
        ["IOControl/gpio_handler.py","GPIO — heartbeat LED, alarm LED, бутон ISR, тахометър RPM"],
        ["Networking/mqtt_client.py","paho-mqtt pub/sub + reconnect + fallback логика"],
        ["Networking/snmp_agent.py","SNMP v2c агент — OID стойности + trap изпращане"],
        ["Storage/data_logger.py","aiosqlite асинхронен логер — readings + events таблици"],
    ]
)

dh2("3.2. React Фронтенд (Vite + React 18)")
dbullets([
    "SensorGauges.jsx — RadialBarChart циферблати за температура и влажност",
    "FanControl.jsx — четири бутона OFF/НИСКА/СРЕДНА/ВИСОКА (disabled в АВТО)",
    "ValveControl.jsx — два toggle превключвателя за клапи",
    "AirflowStatus.jsx — визуален статус на DI-0 и DI-1",
    "ModeToggle.jsx — превключвател АВТО ↔ РЪЧЕН",
    "AlertsPanel.jsx — цветно кодирани известия (info/warning/critical)",
    "HistoricalChart.jsx — LineChart (Recharts) за 1/6/12/24 часа",
    "SimulationPanel.jsx — слайдери + предварителен преглед + бързи сценарии",
    "useWebSocket.js — custom hook с exponential-backoff reconnect",
])

# ── 4. Комуникационни протоколи ───────────────────────────────────────────────
dh1("4. Комуникационни протоколи")
dh2("4.1. MQTT — Mosquitto брокер")
dp("Системата използва MQTT v3.1.1 (paho-mqtt библиотека) за публикуване на телеметрия и получаване на команди от отдалечени клиенти или други Pi устройства.")
dh3("Топик структура")
dtable(
    ["Топик","Тип","Описание"],
    [
        ["datacenter/{node}/status","Pub (retained)","Пълно JSON състояние"],
        ["datacenter/{node}/sensors/temperature","Pub (retained)","Температура float"],
        ["datacenter/{node}/sensors/humidity","Pub (retained)","Влажност float"],
        ["datacenter/{node}/actuators/fan_speed","Pub (retained)","0–3"],
        ["datacenter/{node}/actuators/valve1","Pub (retained)","0/1"],
        ["datacenter/{node}/actuators/valve2","Pub (retained)","0/1"],
        ["datacenter/{node}/airflow/inlet","Pub (retained)","0/1"],
        ["datacenter/{node}/airflow/outlet","Pub (retained)","0/1"],
        ["datacenter/{node}/cmd/fan_speed","Sub","{speed: 0-3}"],
        ["datacenter/{node}/cmd/valve","Sub","{valve_id: 1-2, open: true}"],
        ["datacenter/{node}/cmd/mode","Sub","{mode: auto|manual}"],
        ["datacenter/{node}/cmd/thresholds","Sub","{off: 20, low: 25, medium: 30}"],
    ]
)
dp("Fallback логика: след MQTT_MAX_RECONNECTS (по подразбиране 10) неуспешни опита, системата активира SNMP режим. MQTT опитите продължават на заден план и при успешна връзка fallback се деактивира автоматично.")

dh2("4.2. SNMP v2c агент")
dtable(
    ["OID","Описание","Формат"],
    [
        ["1.3.6.1.4.1.54321.1.1.0","Температура","Integer × 10 (°C)"],
        ["1.3.6.1.4.1.54321.1.2.0","Влажност","Integer × 10 (%)"],
        ["1.3.6.1.4.1.54321.2.1.0","Скорост вентилатор","0=OFF 1=LOW 2=MED 3=HIGH"],
        ["1.3.6.1.4.1.54321.2.2.0","Клапа 1","0=ЗАТВОРЕНА 1=ОТВОРЕНА"],
        ["1.3.6.1.4.1.54321.2.3.0","Клапа 2","0=ЗАТВОРЕНА 1=ОТВОРЕНА"],
        ["1.3.6.1.4.1.54321.3.1.0","Въздух — вход","0/1"],
        ["1.3.6.1.4.1.54321.3.2.0","Въздух — изход","0/1"],
        ["1.3.6.1.4.1.54321.4.1","Trap: висока T°","T ≥ 30 °C"],
        ["1.3.6.1.4.1.54321.4.2","Trap: висока влажност","RH ≥ 70 %"],
    ]
)

# ── 5. Логика на управление ───────────────────────────────────────────────────
dh1("5. Логика на автоматично управление")
dh2("5.1. Температурни прагове")
dtable(
    ["Температура","Вентилатор","Клапа 1 (Рецирк.)","Клапа 2 (Изпуск.)"],
    [
        ["T < 20 °C","ИЗКЛЮЧЕН","ОТВОРЕНА","ЗАТВОРЕНА"],
        ["20 ≤ T < 25 °C","НИСКА (Coil 0)","ОТВОРЕНА","ЗАТВОРЕНА"],
        ["25 ≤ T < 30 °C","СРЕДНА (Coil 1)","ЗАТВОРЕНА","ОТВОРЕНА"],
        ["T ≥ 30 °C","ВИСОКА (Coil 2)","ЗАТВОРЕНА","ОТВОРЕНА"],
    ]
)
dh2("5.2. Fault override — защита при блокирана клапа")
dp("Ако системата е в режим на изпускане (Клапа 2 отворена) но DI-1 (изходящ поток) не е засечен, докато DI-0 (входящ поток) е активен — системата автоматично превключва на рециркулация и записва предупреждение.")

# ── 6. REST API ───────────────────────────────────────────────────────────────
dh1("6. REST API референция")
dtable(
    ["Метод","Endpoint","Тяло / Параметри","Описание"],
    [
        ["GET","/api/status","—","Текущо пълно състояние"],
        ["POST","/api/fan","{speed: 0-3}","Задай скорост (РЪЧЕН режим)"],
        ["POST","/api/valve","{valve_id:1-2, open:bool}","Управлявай клапа"],
        ["POST","/api/mode","{mode: auto|manual}","Смени режим"],
        ["GET /POST","/api/thresholds","{off,low,medium}","Прочети/задай прагове"],
        ["GET","/api/history","?hours=24&limit=500","Исторически данни"],
        ["GET","/api/statistics","—","24ч min/max/avg"],
        ["GET","/api/events","?limit=50","Системни събития"],
        ["GET","/api/snmp","—","SNMP MIB информация"],
        ["GET","/api/mqtt","—","MQTT статус на връзката"],
        ["GET","/api/gpio","—","GPIO пинове и live стойности"],
        ["GET","/api/node","—","Node ID и MQTT топик база"],
        ["GET /POST","/api/simulation","{enabled,temperature,humidity}","Симулационен override"],
        ["WS","/ws","—","WebSocket real-time stream"],
    ]
)

# ── 7. Симулационен режим ─────────────────────────────────────────────────────
dh1("7. Симулационен режим")
dp("Системата поддържа два нива на симулация:")
dh2("7.1. Хардуерна симулация (автоматична)")
dbullets([
    "При липса на BME280/DHT22 — sensor_reader.py генерира реалистични стойности (синусоидален дрейф)",
    "При липса на Modbus RS-485 адаптер — modbus_control.py симулира relay/DI операции",
    "Системата стартира нормално и е напълно функционална без хардуер",
])
dh2("7.2. Ръчен API override (SimulationPanel)")
dp("Чрез React интерфейса или директно чрез API могат да се зададат конкретни стойности за сензора:")
dp("POST /api/simulation  →  {\"enabled\": true, \"temperature\": 31.5, \"humidity\": 68}", bold=True)
dbullets([
    "Слайдер за температура: 10–45 °C с цветова лента (зелен/жълт/червен)",
    "Слайдер за влажност: 10–95 %",
    "Предварителен преглед — показва очакваното поведение ПРЕДИ прилагане",
    "Бързи сценарии: Нормална (22/55) / Топло (27/60) / Горещо (28/63) / Критично (34/71) / Висока влажност (22/82)",
    "Индикатор ⚗ СИМУЛАЦИЯ в заглавието на Sensor панела",
])

# ── 8. Инсталация ─────────────────────────────────────────────────────────────
dh1("8. Инсталация и стартиране")
dh2("8.1. Бързо стартиране (Raspberry Pi)")
dp("bash install.sh", bold=True)
dp("Скриптът автоматично: активира I²C, инсталира Python зависимости, изгражда React, инсталира systemd услуга.")
dp("По желание се инсталира Mosquitto MQTT брокер на същия Pi.")

dh2("8.2. Python зависимости (requirements.txt)")
dbullets([
    "fastapi, uvicorn[standard]",
    "pymodbus >= 3.6",
    "pysnmp >= 6.1",
    "paho-mqtt >= 1.6.1, < 3.0",
    "aiosqlite >= 0.20",
    "wsproto (за WebSocket съвместимост с Python 3.14)",
    "smbus2, adafruit-circuitpython-bme280, RPi.GPIO",
])
dh2("8.3. Ръчно стартиране")
dp("source venv/bin/activate  &&  python main.py", bold=True)
dp("Уеб интерфейс: http://<RPi-IP>:8000  |  API документация: http://<RPi-IP>:8000/docs")

dh2("8.4. systemd услуга")
dp("sudo systemctl enable datacenter_thermal  &&  sudo systemctl start datacenter_thermal", bold=True)

# ── 9. Заключение ─────────────────────────────────────────────────────────────
dh1("9. Заключение")
dp("В рамките на курсовия проект е разработена пълнофункционална разпределена вградена система за термично управление, включваща:")
dbullets([
    "Следене на температура и влажност с SNMP докладване",
    "Управление на вентилатор (3 степени) и 2 клапи чрез Modbus RTU",
    "GPIO управление на индикатори, бутон и тахометър",
    "MQTT телеметрия с multi-node поддръжка и автоматичен SNMP fallback",
    "Уеб базиран React интерфейс с реално-временни данни",
    "Симулационен режим за тестване без физически хардуер",
    "Автоматичен systemd старт и пълна документация",
])
dp("Системата е тествана на Windows с Python 3.14 в симулационен режим и е готова за разполагане на Raspberry Pi 5.")

doc.save(os.path.join(OUT, "dokumentaciya.docx"))
print("✓ dokumentaciya.docx saved")
print("\nAll done! Files in:", OUT)

